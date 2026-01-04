import os
from pptx import Presentation
from pptx.util import Inches, Pt

# ==========================================
# 参数配置区 (Configuration Section)
# ==========================================
# 指定图片所在的目录路径 (请使用绝对路径，例如: r"C:\Users\Data\Images")
# 若留空为 ""，则默认扫描当前脚本所在的目录
INPUT_DIRECTORY = r"D:/数据库复习/24年b" 

# 输出的PPT文件名称
OUTPUT_FILENAME = "D:/数据库复习/24年b/output_presentation.pptx"

# 支持的文件扩展名格式
SUPPORTED_EXTENSIONS = ('.jpg', '.jpeg', '.png', '.bmp', '.gif', '.tiff')
# ==========================================

def get_image_files(directory):
    """
    获取指定目录下所有支持的图片文件，并按文件名排序。
    
    Args:
        directory (str): 目标目录路径。
        
    Returns:
        list: 排序后的图片文件路径列表。
    """
    files = []
    
    # 校验目录是否存在
    if not os.path.exists(directory):
        print(f"错误: 找不到指定的目录 -> {directory}")
        return []

    # 遍历目录查找文件
    for filename in os.listdir(directory):
        # 使用配置区定义的扩展名列表
        if filename.lower().endswith(SUPPORTED_EXTENSIONS):
            files.append(filename)
    
    # 按文件名进行字典序排序
    files.sort()
    return files

def add_image_to_slide(prs, image_path):
    """
    在PPT中新建一页幻灯片，并将图片居中适应放入。
    
    Args:
        prs (Presentation): PPT对象。
        image_path (str): 图片文件路径。
    """
    # 使用空白版式 (索引6通常为空白页)
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # 获取幻灯片的物理尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 插入图片以获取其原始尺寸，初始位置设为0
    # 注意：此时图片可能超出边界，后续步骤将进行调整
    pic = slide.shapes.add_picture(image_path, 0, 0)
    
    # 获取图片的原始宽高
    pic_width = pic.width
    pic_height = pic.height
    
    # 计算宽高比
    image_ratio = pic_width / pic_height
    slide_ratio = slide_width / slide_height
    
    # 逻辑判断：根据宽高比决定缩放基准
    if image_ratio > slide_ratio:
        # 图片较宽，以幻灯片宽度为基准进行缩放
        pic.width = slide_width
        pic.height = int(slide_width / image_ratio)
        # 垂直居中
        pic.top = int((slide_height - pic.height) / 2)
        pic.left = 0
    else:
        # 图片较高，以幻灯片高度为基准进行缩放
        pic.height = slide_height
        pic.width = int(slide_height * image_ratio)
        # 水平居中
        pic.left = int((slide_width - pic.width) / 2)
        pic.top = 0

def main():
    """
    主程序执行入口。
    """
    # 确定目标处理目录
    if INPUT_DIRECTORY:
        target_directory = INPUT_DIRECTORY
    else:
        target_directory = os.getcwd()

    # 确定输出文件的完整路径
    output_path = os.path.join(target_directory, OUTPUT_FILENAME)
    
    print(f"正在初始化...")
    print(f"目标目录: {target_directory}")
    print(f"输出路径: {output_path}")
    
    images = get_image_files(target_directory)
    
    if not images:
        print("未在目标目录下发现支持的图片文件。程序终止。")
        return

    print(f"共发现 {len(images)} 张图片，开始处理...")
    
    # 初始化PPT对象
    prs = Presentation()
    
    for filename in images:
        # 拼接完整的图片路径
        full_image_path = os.path.join(target_directory, filename)
        try:
            add_image_to_slide(prs, full_image_path)
            print(f"已处理: {filename}")
        except Exception as e:
            print(f"处理文件 {filename} 时发生错误: {e}")
            
    try:
        prs.save(output_path)
        print(f"处理完成。文件已保存为: {output_path}")
    except PermissionError:
        print(f"保存失败: 请确保文件 {OUTPUT_FILENAME} 未被其他程序占用。")
    except Exception as e:
        print(f"保存时发生未知错误: {e}")

if __name__ == "__main__":
    main()